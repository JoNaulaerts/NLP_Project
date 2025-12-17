"""
Multi-format Document Uploader for ChromaDB
Supports: PDF, TXT, MD, DOCX, PY, CSV, PPTX
"""
import os
from pathlib import Path
from typing import Dict, Any, List

import chromadb
from langchain_community.document_loaders import (
    PyPDFLoader,
    TextLoader,
    UnstructuredWordDocumentLoader,
    CSVLoader,
)
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

CHROMA_HOST = os.getenv("CHROMA_HOST", "chromadb")
CHROMA_PORT = int(os.getenv("CHROMA_PORT", "8000"))
COLLECTION_NAME = os.getenv("CHROMA_COLLECTION", "ml_materials")


def make_ids(filepath: str, n: int) -> List[str]:
    """Generate unique IDs for document chunks"""
    stem = Path(filepath).stem
    return [f"{stem}_{i}" for i in range(n)]


def load_pptx(filepath: str) -> List[Document]:
    """Extract text from PowerPoint presentations"""
    try:
        from pptx import Presentation
        
        prs = Presentation(filepath)
        documents = []
        path = Path(filepath)
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            # Create document for slide
            if slide_text:
                content = "\n".join(slide_text)
                doc = Document(
                    page_content=content,
                    metadata={
                        "source": path.name,
                        "slide": slide_idx,
                        "file_type": ".pptx"
                    }
                )
                documents.append(doc)
        
        return documents if documents else []
    
    except ImportError:
        raise ImportError("python-pptx not installed. Run: pip install python-pptx")
    except Exception as e:
        raise ValueError(f"Failed to parse PPTX: {str(e)}")


def load_document(filepath: str) -> List[Document]:
    """Load document based on file extension"""
    path = Path(filepath)
    extension = path.suffix.lower()
    
    try:
        if extension == ".pdf":
            loader = PyPDFLoader(filepath)
            return loader.load()
        
        elif extension in [".txt", ".md", ".py"]:
            # Text-based files
            loader = TextLoader(filepath, encoding="utf-8")
            return loader.load()
        
        elif extension == ".docx":
            # Word documents
            loader = UnstructuredWordDocumentLoader(filepath)
            return loader.load()
        
        elif extension == ".pptx":
            # PowerPoint presentations
            return load_pptx(filepath)
        
        elif extension == ".csv":
            # CSV files (each row becomes a document)
            loader = CSVLoader(filepath, encoding="utf-8")
            return loader.load()
        
        else:
            # Fallback: try to read as text
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                content = f.read()
            return [Document(
                page_content=content,
                metadata={"source": path.name, "file_type": extension}
            )]
    
    except Exception as e:
        raise ValueError(f"Failed to load {path.name}: {str(e)}")


def upload_document_to_chromadb(filepath: str) -> Dict[str, Any]:
    """
    Upload any supported document type to ChromaDB
    Supported: .pdf, .txt, .md, .docx, .py, .csv, .pptx
    """
    try:
        path = Path(filepath)
        
        # Load document
        pages = load_document(filepath)
        if not pages:
            return {
                "success": False,
                "message": f"No content in {path.suffix} file.",
                "chunks": 0,
                "pages": 0
            }

        # Split into chunks
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )
        chunks = splitter.split_documents(pages)

        # Prepare data for ChromaDB
        ids = make_ids(filepath, len(chunks))
        documents = [c.page_content for c in chunks]
        metadatas = [
            {
                "source": path.name,
                "file_type": path.suffix,
                "page": c.metadata.get("page", c.metadata.get("slide", 0)),
                "row": c.metadata.get("row", None),  # For CSV
            }
            for c in chunks
        ]

        # Upload to ChromaDB
        client = chromadb.HttpClient(host=CHROMA_HOST, port=CHROMA_PORT)
        collection = client.get_or_create_collection(name=COLLECTION_NAME)
        collection.upsert(ids=ids, documents=documents, metadatas=metadatas)

        return {
            "success": True,
            "message": f"Indexed {path.name} into {COLLECTION_NAME}",
            "chunks": len(chunks),
            "pages": len(pages),
        }
    
    except Exception as e:
        return {
            "success": False,
            "message": f"Error: {str(e)}",
            "chunks": 0,
            "pages": 0
        }


# Backward compatibility aliases
upload_pdf_to_chromadb = upload_document_to_chromadb
